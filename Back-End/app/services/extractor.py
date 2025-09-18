import re, unicodedata
from pathlib import Path
from typing import List
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

SLIDE_BREAK = "\n\n---SLIDE_BREAK---\n\n"

_DASH_MAP = {
    "\u00AD": "",  "\u2010": "-", "\u2011": "-", "\u2012": "-",
    "\u2013": "–", "\u2014": "—", "\u2212": "-", "\u2043": "-"
}
_TRANS = str.maketrans(_DASH_MAP)
_BULLETS = {"-", "•", "–", "—", "∙", "‣"}

def _normalize_block(t: str) -> str:
    t = unicodedata.normalize("NFKC", t or "").translate(_TRANS)
    t = re.sub(r'(?<=\w)-\s*\n\s*(?=\w)', "", t)        
    t = re.sub(r'[ \t]+\n', '\n', t)
    t = re.sub(r'\n{3,}', '\n\n', t)
    return t

def _reflow_lines_keep_bullets(t: str) -> str:
    """
    Junta palavras quebradas linha a linha em parágrafos.
    Linha vazia => quebra de parágrafo.
    Linha que é só o marcador de bullet começa um item novo.
    """
    lines = [ln.strip() for ln in t.splitlines()]
    out: List[str] = []
    buf: List[str] = []
    bullet_mode = False

    def flush():
        nonlocal buf, bullet_mode
        if not buf:
            return
        joined = " ".join(buf).strip()
        if bullet_mode:
            out.append(f"- {joined}")
        else:
            out.append(joined)
        buf = []
        bullet_mode = False

    for ln in lines:
        if ln == "":
            flush()
            out.append("")  
            continue
        if ln in _BULLETS:       
            flush()
            bullet_mode = True
            continue
        buf.append(ln)

    flush()
    
    result: List[str] = []
    prev_blank = False
    for s in out:
        if s == "":
            if not prev_blank:
                result.append("")
            prev_blank = True
        else:
            result.append(s)
            prev_blank = False
    return "\n".join(result).strip()

def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".docx":
        doc = Document(file_path)
        raw = "\n".join((p.text or "").replace("\r", "") for p in doc.paragraphs)
        return _reflow_lines_keep_bullets(_normalize_block(raw))

    if ext == ".pptx":
        prs = Presentation(file_path)
        slides_out: List[str] = []
        for slide in prs.slides:
            lines: List[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join((run.text or "") for run in para.runs) or (para.text or "")
                        txt = txt.replace("\r", "\n")
                        if txt.strip():
                            lines.append(txt)
            slide_text = _reflow_lines_keep_bullets(_normalize_block("\n".join(lines)))
            slides_out.append(slide_text)
        return SLIDE_BREAK.join(slides_out)

    if ext == ".pdf":
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            txt = page.extract_text() or ""
            pages.append(_reflow_lines_keep_bullets(_normalize_block(txt)))
        return "\n\n".join(pages)

    if ext == ".ppt":
        raise ValueError("Formato .ppt não suportado. Converta para .pptx.")
    raise ValueError("Formato de arquivo não suportado. Use .docx, .pptx ou .pdf.")

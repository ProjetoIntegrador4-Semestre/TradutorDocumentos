import os
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
import re


SLIDE_BREAK = "\n\n---SLIDE_BREAK---\n\n"


def generate_file(original_path: str, translated_text: str) -> str:
    ext = os.path.splitext(original_path)[1].lower()
    base = os.path.basename(original_path)
    out_path = f"temp/translated_{base}"

    if ext == ".docx":
        doc = Document()
        for line in translated_text.split("\n"):
            doc.add_paragraph(line)
        doc.save(out_path)
        return out_path

    if ext == ".pptx":
        prs = Presentation()
        chunks = [c.strip() for c in translated_text.split(SLIDE_BREAK)]
        chunks = [c for c in chunks if c != ""] or [""]

        for chunk in chunks:
            lines = [l for l in chunk.split("\n") if l.strip()]

            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Título + Conteúdo
            title_placeholder = slide.shapes.title
            content_placeholder = slide.placeholders[1].text_frame

            # título = primeira linha (ou "Slide")
            title_placeholder.text = (lines[0] if lines else "Slide")[:120]

            # conteúdo: limpar parágrafo default e escrever linha a linha
            content_placeholder.clear()
            if len(lines) > 1:
                for i, line in enumerate(lines[1:]):
                    if i == 0:
                        p = content_placeholder.paragraphs[0]
                    else:
                        p = content_placeholder.add_paragraph()
                    p.text = line
                    p.level = 0

        prs.save(out_path)
        return out_path

    if ext == ".pdf":
        c = canvas.Canvas(out_path, pagesize=A4)
        t = c.beginText(50, 800)
        for line in translated_text.split("\n"):
            t.textLine(line)
        c.drawText(t)
        c.save()
        return out_path

    raise ValueError("Formato de arquivo não suportado.")

